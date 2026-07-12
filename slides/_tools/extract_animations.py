"""
Extract animation sequences from PPTX files.

For each deck, produces slides/XX/extracted/animations.json describing:
- Which shapes appear on each click (click groups)
- Shape types (title, text, image, table, chart)
- Text content and image filenames

This data drives v-click reconstruction in Slidev.
"""

import json
import re
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
import lxml.etree as etree

PPTX_NS = {
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
}

SLIDES_DIR = Path(__file__).parent.parent


def get_shape_text(shape) -> str:
    """Extract text content from a shape."""
    try:
        if shape.has_text_frame:
            lines = []
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    lines.append(text)
            return '\n'.join(lines)
    except Exception:
        pass
    return ''


def get_shape_type(shape) -> str:
    """Classify a shape as title/text/image/table/chart/other."""
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
        # Placeholder types
        if shape.is_placeholder:
            ph = shape.placeholder_format
            if ph.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                return 'title'
            if ph.type == PP_PLACEHOLDER.SUBTITLE:
                return 'subtitle'
            if ph.type == PP_PLACEHOLDER.PICTURE:
                return 'image'
            if ph.type == PP_PLACEHOLDER.BODY:
                return 'text'
        # Shape types
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return 'image'
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            return 'table'
        if shape.shape_type == MSO_SHAPE_TYPE.CHART:
            return 'chart'
        if shape.has_text_frame:
            return 'text'
    except Exception:
        pass
    return 'other'


def get_image_filename(shape, slide_images: dict) -> str:
    """Try to find the image filename for a picture shape."""
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            img = shape.image
            # Match by content hash against known images
            img_hash = hash(img.blob)
            if img_hash in slide_images:
                return slide_images[img_hash]
    except Exception:
        pass
    return ''


def build_shape_index(prs: Presentation) -> dict:
    """Build a map of shape ID → shape object across all slides."""
    index = {}
    for slide in prs.slides:
        for shape in slide.shapes:
            index[shape.shape_id] = shape
    return index


def extract_click_groups(slide_elem, shapes_on_slide: dict) -> dict:
    """
    Parse <p:timing> to extract click-triggered animation groups.

    Each clickEffect par in OOXML represents one user click.
    Paragraph-level animations use <p:pRg st="N" end="N"/> to specify
    which paragraph in a text box becomes visible.

    Returns dict: {
        (shape_id, para_idx_or_None): click_index,
        ...
    }
    For whole-shape animations, para_idx is None.
    For paragraph animations, para_idx is the paragraph index.
    """
    ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    timing = slide_elem.find(f'.//{{{ns}}}timing')
    if timing is None:
        return {}

    animated = {}  # (shape_id, para_idx) -> click_index
    click_index = 0

    for par in timing.findall(f'.//{{{ns}}}par'):
        cTn = par.find(f'{{{ns}}}cTn')
        if cTn is None:
            continue
        if cTn.get('nodeType') != 'clickEffect':
            continue
        click_index += 1

        # Find all spTgt in this clickEffect subtree
        sp_tgts = par.findall(f'.//{{{ns}}}spTgt')
        for sp_tgt in sp_tgts:
            spid = sp_tgt.get('spid')
            if not spid:
                continue
            shape_id = int(spid)

            # Check for paragraph-level targeting via <p:txEl><p:pRg st="N" end="N"/>
            tx_el = sp_tgt.find(f'{{{ns}}}txEl')
            if tx_el is not None:
                p_rg = tx_el.find(f'{{{ns}}}pRg')
                if p_rg is not None:
                    st = int(p_rg.get('st', 0))
                    end = int(p_rg.get('end', st))
                    for para_idx in range(st, end + 1):
                        key = (shape_id, para_idx)
                        if key not in animated:
                            animated[key] = click_index
                    continue
            # Whole-shape animation
            key = (shape_id, None)
            if key not in animated:
                animated[key] = click_index

    return animated


def extract_slide_animations(slide, slide_num: int) -> dict:
    """Extract animation data for a single slide."""
    # animated: {(shape_id, para_idx_or_None): click_index}
    animated = extract_click_groups(slide.element, {})

    max_click = max(animated.values(), default=0)

    # Build items list: each item is a text paragraph or whole shape revealed at a click
    items = []
    for shape in slide.shapes:
        shape_type = get_shape_type(shape)

        # Whole-shape animation?
        whole_key = (shape.shape_id, None)
        if whole_key in animated:
            items.append({
                'shape_id': shape.shape_id,
                'shape_name': shape.name,
                'type': shape_type,
                'click': animated[whole_key],
                'text': get_shape_text(shape)[:200] if shape.has_text_frame else '',
                'is_image': shape_type == 'image',
            })
        elif shape.has_text_frame:
            # Check paragraph-level animations
            # Sub-paragraphs (level > 0) inherit the click of their parent level-0 paragraph
            last_parent_click = 0
            for para_idx, para in enumerate(shape.text_frame.paragraphs):
                para_key = (shape.shape_id, para_idx)
                text = para.text.strip()
                level = para.level  # 0 = top level, 1+ = sub-levels

                if para_key in animated:
                    click = animated[para_key]
                    if level == 0:
                        last_parent_click = click
                elif level > 0 and last_parent_click > 0:
                    # Sub-paragraph: inherit parent's click
                    click = last_parent_click
                else:
                    click = 0  # visible from start

                if text:
                    items.append({
                        'shape_id': shape.shape_id,
                        'shape_name': shape.name,
                        'type': 'paragraph',
                        'click': click,
                        'para_idx': para_idx,
                        'level': level,
                        'text': text[:200],
                        'is_image': False,
                    })
        elif shape_type == 'image' and whole_key not in animated:
            # Image visible from start
            items.append({
                'shape_id': shape.shape_id,
                'shape_name': shape.name,
                'type': 'image',
                'click': 0,
                'text': '',
                'is_image': True,
            })

    # Group by click index
    click_groups = {}
    for item in items:
        c = item['click']
        if c not in click_groups:
            click_groups[c] = []
        click_groups[c].append(item)

    return {
        'number': slide_num,
        'has_animations': len(animated) > 0,
        'animation_count': len(animated),
        'click_count': max_click,
        'click_groups': [
            {'click': k, 'items': click_groups[k]}
            for k in sorted(click_groups.keys())
        ]
    }


def process_deck(deck_dir: Path) -> dict:
    """Process one PPTX deck and return animation data."""
    original_dir = deck_dir / 'original'
    pptx_files = list(original_dir.glob('*.pptx'))
    if not pptx_files:
        return None

    pptx_path = pptx_files[0]
    print(f"  Processing: {pptx_path.name}")

    prs = Presentation(pptx_path)

    slides_data = []
    animated_count = 0
    for i, slide in enumerate(prs.slides):
        slide_data = extract_slide_animations(slide, i + 1)
        slides_data.append(slide_data)
        if slide_data['has_animations']:
            animated_count += 1

    return {
        'filename': pptx_path.name,
        'slide_count': len(prs.slides),
        'animated_slide_count': animated_count,
        'slides': slides_data,
    }


def main():
    decks = sorted([d for d in SLIDES_DIR.iterdir()
                    if d.is_dir() and (d / 'original').exists()])

    target = sys.argv[1] if len(sys.argv) > 1 else None

    for deck_dir in decks:
        if target and deck_dir.name != target:
            continue

        print(f"\n[{deck_dir.name}]")
        data = process_deck(deck_dir)
        if data is None:
            print("  No PPTX found, skipping.")
            continue

        output_path = deck_dir / 'extracted' / 'animations.json'
        output_path.parent.mkdir(parents=True, exist_ok=True)
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(data, f, ensure_ascii=False, indent=2)

        print(f"  Slides: {data['slide_count']}, Animated: {data['animated_slide_count']}")
        print(f"  -> {output_path}")


if __name__ == '__main__':
    main()
