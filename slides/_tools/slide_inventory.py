"""
slide_inventory.py - Batch inventory of student presentations.

Scans multiple source directories for presentation files (.pptx, .ppt, .pdf, .odp)
and generates a structured catalog with metadata, topic classification, and quality hints.

Usage:
    python slide_inventory.py scan [--sources FILE] [--output JSON]
    python slide_inventory.py classify [--catalog JSON]
    python slide_inventory.py report [--catalog JSON]
"""

import argparse
import json
import os
import sys
import time
from dataclasses import dataclass, field, asdict
from pathlib import Path
from typing import Optional


# ---------------------------------------------------------------------------
# Topic classification keywords
# ---------------------------------------------------------------------------

TOPIC_KEYWORDS = {
    "search-csp": [
        "sudoku", "backtracking", "constraint", "csp", "coloration", "graphe",
        "dancing link", "dlx", "norvig", "a*", "bfs", "dfs", "genetic",
        "exploration", "heuristique", "optimisation", "knapsack", "tsp",
    ],
    "game-theory": [
        "game", "jeu", "nash", "minimax", "alpha-beta", "go ", "poker",
        "trueskill", "cooperative", "strateg", "equilibr", "theorie des jeux",
        "cfr", "regret", "echec", "chess",
    ],
    "ml-dl": [
        "machine learning", "apprentissage", "neural", "deep learning", "cnn",
        "classification", "regression", "random forest", "svm", "tensorflow",
        "sentiment", "cancer", "detection", "image", "transfer learning",
        "reseau de neurone", "perceptron", "pytorch", "keras", "sklearn",
    ],
    "probas": [
        "probabilit", "bayes", "markov", "hmm", "infer", "gaussian",
        "incertitude", "monte carlo", "stochast",
    ],
    "nlp": [
        "langage", "nlp", "chatbot", "sentiment", "text", "bert", "gpt",
        "traitement.*langue", "natural language", "tok",
    ],
    "trading": [
        "trading", "finance", "portfolio", "crypto", "bitcoin", "quantconnect",
        "algorithmique.*trading", "bourse", "action",
    ],
    "genai": [
        "generati", "gan", "diffusion", "dall-e", "stable diffusion",
        "ia generative", "llm", "chatgpt", "image.*generat",
    ],
    "medical": [
        "medical", "cancer", "asthme", "diagnostic", "sante", "health",
        "tumeur", "breast", "clinical",
    ],
}

# Presentation file extensions to scan
PRESENTATION_EXTENSIONS = {".pptx", ".ppt", ".odp", ".pdf", ".gslides", ".key"}


@dataclass
class PresentationEntry:
    """Metadata for a single presentation file."""
    filepath: str
    filename: str
    extension: str
    file_size_mb: float
    source_dir: str
    school: str = ""
    year: str = ""
    topic: str = "misc"
    topic_confidence: float = 0.0
    slide_count: int = 0
    image_count: int = 0
    title_text: str = ""


def scan_directory(root_dir: str, source_label: str = "") -> list:
    """Scan a directory tree for presentation files.

    Args:
        root_dir: Root directory to scan.
        source_label: Label for this source (e.g., "dropbox", "gdrive").

    Returns:
        List of PresentationEntry objects.
    """
    entries = []
    root = Path(root_dir)

    if not root.exists():
        print(f"WARNING: Directory not found: {root_dir}")
        return entries

    for filepath in root.rglob("*"):
        if filepath.suffix.lower() in PRESENTATION_EXTENSIONS:
            try:
                size_mb = filepath.stat().st_size / (1024 * 1024)
            except OSError:
                size_mb = 0

            entry = PresentationEntry(
                filepath=str(filepath),
                filename=filepath.name,
                extension=filepath.suffix.lower(),
                file_size_mb=round(size_mb, 2),
                source_dir=source_label or root_dir,
            )

            # Try to extract school and year from path
            parts = filepath.relative_to(root).parts
            for part in parts:
                if part in ("EPF", "ECE", "ORT", "Epita", "ESGF", "ESIEE",
                            "Simplon", "Aston", "Monaco"):
                    entry.school = part
                if part.isdigit() and len(part) == 4 and 2015 <= int(part) <= 2030:
                    entry.year = part

            entries.append(entry)

    return entries


def classify_entry(entry: PresentationEntry) -> PresentationEntry:
    """Classify a presentation entry by topic using keyword matching.

    Args:
        entry: PresentationEntry to classify.

    Returns:
        The same entry with topic and topic_confidence set.
    """
    # Build search text from filename + path
    search_text = (entry.filename + " " + entry.filepath).lower()

    best_topic = "misc"
    best_score = 0

    for topic, keywords in TOPIC_KEYWORDS.items():
        score = 0
        for keyword in keywords:
            if keyword.lower() in search_text:
                score += 1
        if score > best_score:
            best_score = score
            best_topic = topic

    entry.topic = best_topic
    entry.topic_confidence = min(best_score / 3.0, 1.0)  # Normalize to [0, 1]
    return entry


def enrich_pptx_metadata(entry: PresentationEntry) -> PresentationEntry:
    """Enrich a PPTX entry with slide count and title (via python-pptx).

    Args:
        entry: PresentationEntry for a .pptx file.

    Returns:
        The same entry with slide_count, image_count, title_text set.
    """
    if entry.extension != ".pptx":
        return entry

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        prs = Presentation(entry.filepath)
        entry.slide_count = len(prs.slides)

        # Get title from first slide
        if prs.slides and prs.slides[0].shapes.title:
            entry.title_text = prs.slides[0].shapes.title.text.strip()[:200]

        # Count images
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    entry.image_count += 1

    except Exception as e:
        print(f"  WARNING: Could not read {entry.filename}: {e}")

    return entry


def build_catalog(sources: list, enrich: bool = True) -> list:
    """Build a full catalog from multiple source directories.

    Args:
        sources: List of (directory_path, source_label) tuples.
        enrich: Whether to enrich PPTX files with metadata (slower).

    Returns:
        List of PresentationEntry objects.
    """
    all_entries = []

    for directory, label in sources:
        print(f"Scanning: {label} ({directory})")
        entries = scan_directory(directory, label)
        print(f"  Found {len(entries)} presentation files")
        all_entries.extend(entries)

    print(f"\nTotal files found: {len(all_entries)}")

    # Classify all entries
    print("Classifying by topic...")
    for entry in all_entries:
        classify_entry(entry)

    # Enrich PPTX files
    if enrich:
        pptx_entries = [e for e in all_entries if e.extension == ".pptx"]
        print(f"Enriching {len(pptx_entries)} PPTX files with metadata...")
        for i, entry in enumerate(pptx_entries, 1):
            enrich_pptx_metadata(entry)
            if i % 20 == 0:
                print(f"  Enriched {i}/{len(pptx_entries)}")

    return all_entries


def save_catalog(entries: list, output_path: str):
    """Save catalog to JSON file."""
    data = [asdict(e) for e in entries]
    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=2, ensure_ascii=False)
    print(f"Catalog saved to: {output_path} ({len(entries)} entries)")


def print_report(entries: list):
    """Print a summary report of the catalog."""
    print(f"\n{'='*60}")
    print(f"PRESENTATION CATALOG REPORT")
    print(f"{'='*60}")
    print(f"Total entries: {len(entries)}")

    # By extension
    ext_counts = {}
    for e in entries:
        ext_counts[e.extension] = ext_counts.get(e.extension, 0) + 1
    print(f"\nBy format:")
    for ext, count in sorted(ext_counts.items(), key=lambda x: -x[1]):
        print(f"  {ext}: {count}")

    # By school
    school_counts = {}
    for e in entries:
        school = e.school or "(unknown)"
        school_counts[school] = school_counts.get(school, 0) + 1
    print(f"\nBy school:")
    for school, count in sorted(school_counts.items(), key=lambda x: -x[1]):
        print(f"  {school}: {count}")

    # By topic
    topic_counts = {}
    for e in entries:
        topic_counts[e.topic] = topic_counts.get(e.topic, 0) + 1
    print(f"\nBy topic:")
    for topic, count in sorted(topic_counts.items(), key=lambda x: -x[1]):
        print(f"  {topic}: {count}")

    # By source
    source_counts = {}
    for e in entries:
        source_counts[e.source_dir] = source_counts.get(e.source_dir, 0) + 1
    print(f"\nBy source:")
    for source, count in sorted(source_counts.items(), key=lambda x: -x[1]):
        print(f"  {source}: {count}")

    # Top PPTX by visual richness (image_count / slide_count)
    pptx_entries = [e for e in entries if e.extension == ".pptx" and e.slide_count > 0]
    if pptx_entries:
        pptx_entries.sort(key=lambda e: e.image_count / max(e.slide_count, 1), reverse=True)
        print(f"\nTop 20 visually richest PPTX:")
        for e in pptx_entries[:20]:
            density = e.image_count / e.slide_count
            print(f"  [{e.topic}] {e.filename} ({e.slide_count}s, {e.image_count}img, density={density:.1f}) - {e.school}")


# ---------------------------------------------------------------------------
# Default sources
# ---------------------------------------------------------------------------

DEFAULT_SOURCES = [
    (r"D:\Dropbox\IA101\2017", "dropbox-2017"),
    (r"D:\Dropbox\IA101\2018", "dropbox-2018"),
    (r"D:\Dropbox\IA101\2019", "dropbox-2019"),
    (r"D:\Dropbox\IA101\2020", "dropbox-2020"),
    (r"D:\Dropbox\IA101\2021", "dropbox-2021"),
    (r"D:\Dropbox\IA101\Vulgarisation", "dropbox-vulgarisation"),
    (r"G:\Mon Drive\MyIA\Formation\EPF", "gdrive-EPF"),
    (r"G:\Mon Drive\MyIA\Formation\ECE", "gdrive-ECE"),
    (r"G:\Mon Drive\MyIA\Formation\Epita", "gdrive-Epita"),
    (r"G:\Mon Drive\MyIA\Formation\ESGF", "gdrive-ESGF"),
]


def main():
    parser = argparse.ArgumentParser(description="Student presentation inventory tool")
    subparsers = parser.add_subparsers(dest="command")

    p_scan = subparsers.add_parser("scan", help="Scan directories and build catalog")
    p_scan.add_argument("--output", default="d:/dev/CoursIA/slides/_assets/student-catalog.json",
                        help="Output JSON path")
    p_scan.add_argument("--no-enrich", action="store_true",
                        help="Skip PPTX metadata enrichment (faster)")

    p_report = subparsers.add_parser("report", help="Print report from existing catalog")
    p_report.add_argument("--catalog", default="d:/dev/CoursIA/slides/_assets/student-catalog.json",
                          help="Catalog JSON path")

    args = parser.parse_args()

    if not args.command:
        parser.print_help()
        return

    if args.command == "scan":
        entries = build_catalog(DEFAULT_SOURCES, enrich=not args.no_enrich)
        save_catalog(entries, args.output)
        print_report(entries)

    elif args.command == "report":
        with open(args.catalog, "r", encoding="utf-8") as f:
            data = json.load(f)
        entries = [PresentationEntry(**d) for d in data]
        print_report(entries)


if __name__ == "__main__":
    main()
