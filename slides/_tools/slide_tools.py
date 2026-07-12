"""
slide_tools.py - Outils d'extraction, de rendu et d'analyse de slides PowerPoint et Marp.

Usage:
    python slide_tools.py render <pptx_path> [--outdir DIR] [--width W] [--height H]
    python slide_tools.py extract <pptx_path> [--outdir DIR]
    python slide_tools.py inventory <pptx_path> [--output JSON_PATH]
    python slide_tools.py batch-render <directory> [--outdir DIR]
    python slide_tools.py batch-extract <directory> [--outdir DIR]
    python slide_tools.py marp-render <deck_path> [--scale N]
    python slide_tools.py marp-render-all [--scale N]
    python slide_tools.py check-refs <deck_path>
    python slide_tools.py compare <deck_path>

Requires:
    - pywin32 (COM automation for PowerPoint rendering)
    - python-pptx (PPTX reading/writing)
    - markitdown (text extraction to markdown)
    - marp-cli (npm: @marp-team/marp-cli, for Marp rendering)
"""

import argparse
import json
import os
import sys
import time
from dataclasses import dataclass, field, asdict
from pathlib import Path
from typing import Optional


# ---------------------------------------------------------------------------
# Data classes
# ---------------------------------------------------------------------------

@dataclass
class SlideInfo:
    """Metadata for a single slide."""
    number: int
    title: str = ""
    text_content: str = ""
    word_count: int = 0
    image_count: int = 0
    shape_count: int = 0
    has_title: bool = False
    layout_name: str = ""
    notes: str = ""


@dataclass
class DeckInventory:
    """Full inventory of a PPTX deck."""
    filename: str
    filepath: str
    slide_count: int = 0
    total_images: int = 0
    total_words: int = 0
    file_size_mb: float = 0.0
    slides: list = field(default_factory=list)


# ---------------------------------------------------------------------------
# Rendering (PowerPoint COM via pywin32)
# ---------------------------------------------------------------------------

def render_slides(pptx_path: str, output_dir: str, width: int = 1920, height: int = 1080) -> int:
    """Export each slide of a PPTX to PNG via PowerPoint COM automation.

    Args:
        pptx_path: Path to the PPTX file.
        output_dir: Directory where PNGs will be saved.
        width: Output image width in pixels.
        height: Output image height in pixels.

    Returns:
        Number of slides rendered.
    """
    try:
        import win32com.client
        import pythoncom
    except ImportError:
        print("ERROR: pywin32 is required for rendering. Install with: pip install pywin32")
        sys.exit(1)

    # Ensure absolute paths (COM requires them)
    pptx_path = os.path.abspath(pptx_path)
    output_dir = os.path.abspath(output_dir)

    if not os.path.exists(pptx_path):
        print(f"ERROR: File not found: {pptx_path}")
        sys.exit(1)

    os.makedirs(output_dir, exist_ok=True)

    # Initialize COM in this thread
    pythoncom.CoInitialize()

    try:
        ppt = win32com.client.Dispatch("PowerPoint.Application")
        # Open presentation read-only, no window
        presentation = ppt.Presentations.Open(
            pptx_path,
            ReadOnly=True,
            Untitled=False,
            WithWindow=False
        )

        slide_count = presentation.Slides.Count
        print(f"Rendering {slide_count} slides from: {os.path.basename(pptx_path)}")

        for i in range(1, slide_count + 1):
            slide = presentation.Slides(i)
            out_path = os.path.join(output_dir, f"slide_{i:02d}.png")
            slide.Export(out_path, "PNG", width, height)
            if i % 10 == 0 or i == slide_count:
                print(f"  Rendered {i}/{slide_count}")

        presentation.Close()
        print(f"Done: {slide_count} slides -> {output_dir}")
        return slide_count

    except Exception as e:
        print(f"ERROR during rendering: {e}")
        raise
    finally:
        pythoncom.CoUninitialize()


def batch_render(directory: str, output_base: str = None, width: int = 1920, height: int = 1080) -> dict:
    """Render all PPTX files in a directory.

    Args:
        directory: Directory containing PPTX files.
        output_base: Base output directory. If None, renders to extracted/renders/ subdirs.
        width: Output image width.
        height: Output image height.

    Returns:
        Dict mapping filename -> slide count.
    """
    results = {}
    pptx_files = sorted(Path(directory).rglob("*.pptx"))

    if not pptx_files:
        print(f"No PPTX files found in: {directory}")
        return results

    print(f"Found {len(pptx_files)} PPTX files to render")

    for pptx_path in pptx_files:
        if output_base:
            out_dir = os.path.join(output_base, pptx_path.stem)
        else:
            out_dir = os.path.join(os.path.dirname(pptx_path), "renders")

        try:
            count = render_slides(str(pptx_path), out_dir, width, height)
            results[pptx_path.name] = count
        except Exception as e:
            print(f"FAILED: {pptx_path.name} - {e}")
            results[pptx_path.name] = -1

    return results


# ---------------------------------------------------------------------------
# Marp rendering (marp-cli)
# ---------------------------------------------------------------------------

# Root of the slides directory (relative to this script)
SLIDES_ROOT = Path(__file__).resolve().parent.parent
THEME_PATH = SLIDES_ROOT / "themes" / "ia101.css"

# All deck directory names
DECK_NAMES = [
    "01-introduction", "02-resolution-problemes", "03-logique",
    "04-probabilites", "05-theorie-des-jeux", "06-apprentissage",
    "07-elargissements", "08-ia-generative",
    "S1-argumentation", "S2-ia-exploratoire-symbolique",
    "S3-acculturation", "S4-trading-algorithmique",
]


def render_marp(deck_path: str, scale: int = 2) -> list:
    """Render a Marp slides.md to individual PNG images via marp-cli.

    Args:
        deck_path: Path to the deck directory (must contain slides.md).
        scale: Image scale factor (default 2 for high-DPI).

    Returns:
        List of generated PNG file paths.
    """
    import subprocess

    deck_path = Path(deck_path).resolve()
    slides_md = deck_path / "slides.md"
    output_dir = deck_path / "output" / "marp_renders"

    if not slides_md.exists():
        print(f"ERROR: {slides_md} not found")
        return []

    os.makedirs(output_dir, exist_ok=True)

    # Output base: slide.png -> produces slide.001.png, slide.002.png, ...
    output_base = output_dir / "slide.png"

    cmd = [
        "npx", "@marp-team/marp-cli",
        str(slides_md),
        "--images", "png",
        "--image-scale", str(scale),
        "--html",
        "--allow-local-files",
        "--theme-set", str(THEME_PATH),
        "-o", str(output_base),
    ]

    print(f"Rendering Marp: {deck_path.name}")
    try:
        result = subprocess.run(
            cmd, capture_output=True, text=True, timeout=120,
            cwd=str(SLIDES_ROOT),
        )
        if result.returncode != 0:
            print(f"  ERROR: marp-cli failed:\n{result.stderr}")
            return []
    except subprocess.TimeoutExpired:
        print("  ERROR: marp-cli timed out (120s)")
        return []
    except FileNotFoundError:
        print("  ERROR: marp-cli not found. Install with: npm i -g @marp-team/marp-cli")
        return []

    # Collect generated PNGs
    pngs = sorted(output_dir.glob("slide.*.png"))
    print(f"  Generated {len(pngs)} slides -> {output_dir}")
    return [str(p) for p in pngs]


def render_all_marp(scale: int = 2) -> dict:
    """Render all Marp decks to PNG.

    Returns:
        Dict mapping deck name -> number of slides rendered.
    """
    results = {}
    for deck_name in DECK_NAMES:
        deck_path = SLIDES_ROOT / deck_name
        if (deck_path / "slides.md").exists():
            pngs = render_marp(str(deck_path), scale)
            results[deck_name] = len(pngs)
        else:
            print(f"  SKIP: {deck_name} (no slides.md)")
            results[deck_name] = 0
    return results


def check_image_refs(deck_path: str) -> dict:
    """Check that all image references in slides.md exist on disk.

    Args:
        deck_path: Path to the deck directory.

    Returns:
        Dict with 'found', 'missing', and 'unreferenced' image lists.
    """
    import re

    deck_path = Path(deck_path).resolve()
    slides_md = deck_path / "slides.md"
    images_dir = deck_path / "images"

    if not slides_md.exists():
        print(f"ERROR: {slides_md} not found")
        return {}

    content = slides_md.read_text(encoding="utf-8")

    # Find all image references in markdown and HTML
    md_refs = set(re.findall(r'images/[^\s\)\]"]+', content))

    # Check existence
    found = []
    missing = []
    for ref in sorted(md_refs):
        full_path = deck_path / ref
        if full_path.exists():
            found.append(ref)
        else:
            missing.append(ref)

    # Find unreferenced images
    unreferenced = []
    if images_dir.exists():
        on_disk = {f"images/{f.name}" for f in images_dir.iterdir() if f.is_file()}
        unreferenced = sorted(on_disk - md_refs)

    result = {
        "found": found,
        "missing": missing,
        "unreferenced": unreferenced,
        "total_refs": len(md_refs),
        "total_on_disk": len(found) + len(unreferenced),
    }

    print(f"Image refs in {deck_path.name}:")
    print(f"  Referenced: {len(md_refs)} ({len(found)} found, {len(missing)} missing)")
    if missing:
        for m in missing:
            print(f"    MISSING: {m}")
    if unreferenced:
        print(f"  Unreferenced on disk: {len(unreferenced)}")

    return result


def compare_slide_counts(deck_path: str) -> dict:
    """Compare PPTX vs Marp slide counts for a deck.

    Args:
        deck_path: Path to the deck directory.

    Returns:
        Dict with pptx_count, marp_count, and drift.
    """
    deck_path = Path(deck_path).resolve()

    # Count PPTX renders
    pptx_renders = sorted((deck_path / "extracted" / "renders").glob("slide_*.png"))
    pptx_count = len(pptx_renders)

    # Count Marp slides (by --- separators)
    slides_md = deck_path / "slides.md"
    marp_count = 0
    if slides_md.exists():
        content = slides_md.read_text(encoding="utf-8")
        # Count slide separators (--- at start of line, not in frontmatter)
        lines = content.split("\n")
        in_frontmatter = False
        for i, line in enumerate(lines):
            if i == 0 and line.strip() == "---":
                in_frontmatter = True
                continue
            if in_frontmatter and line.strip() == "---":
                in_frontmatter = False
                continue
            if not in_frontmatter and line.strip() == "---":
                marp_count += 1
        marp_count += 1  # First slide (before any ---)

    result = {
        "deck": deck_path.name,
        "pptx_slides": pptx_count,
        "marp_slides": marp_count,
        "drift": marp_count - pptx_count,
    }

    drift_str = f"+{result['drift']}" if result['drift'] > 0 else str(result['drift'])
    print(f"{deck_path.name}: PPTX={pptx_count}, Marp={marp_count}, drift={drift_str}")
    return result


# ---------------------------------------------------------------------------
# Text extraction (MarkItDown)
# ---------------------------------------------------------------------------

def extract_text(pptx_path: str, output_path: str = None) -> str:
    """Extract text content from PPTX using MarkItDown.

    Args:
        pptx_path: Path to the PPTX file.
        output_path: Optional path to save the markdown output.

    Returns:
        Markdown text content.
    """
    try:
        from markitdown import MarkItDown
    except ImportError:
        print("ERROR: markitdown is required. Install with: pip install markitdown")
        sys.exit(1)

    md = MarkItDown()
    result = md.convert(pptx_path)
    text = result.text_content

    if output_path:
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(text)
        print(f"Text extracted to: {output_path}")

    return text


# ---------------------------------------------------------------------------
# Image extraction (python-pptx)
# ---------------------------------------------------------------------------

def extract_images(pptx_path: str, output_dir: str) -> int:
    """Extract all embedded images from a PPTX file.

    Args:
        pptx_path: Path to the PPTX file.
        output_dir: Directory to save extracted images.

    Returns:
        Number of images extracted.
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(pptx_path)
    os.makedirs(output_dir, exist_ok=True)
    count = 0

    for slide_num, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                ext = image.content_type.split("/")[-1]
                if ext == "jpeg":
                    ext = "jpg"
                elif ext == "x-emf":
                    ext = "emf"
                elif ext == "x-wmf":
                    ext = "wmf"
                filename = f"slide_{slide_num:02d}_img_{count:03d}.{ext}"
                filepath = os.path.join(output_dir, filename)
                with open(filepath, "wb") as f:
                    f.write(image.blob)
                count += 1

    print(f"Extracted {count} images from {os.path.basename(pptx_path)} -> {output_dir}")
    return count


# ---------------------------------------------------------------------------
# Inventory generation (python-pptx)
# ---------------------------------------------------------------------------

def generate_inventory(pptx_path: str, output_path: str = None) -> DeckInventory:
    """Generate a structured inventory of a PPTX deck.

    Args:
        pptx_path: Path to the PPTX file.
        output_path: Optional path to save the JSON inventory.

    Returns:
        DeckInventory object.
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(pptx_path)
    file_size = os.path.getsize(pptx_path) / (1024 * 1024)

    inventory = DeckInventory(
        filename=os.path.basename(pptx_path),
        filepath=os.path.abspath(pptx_path),
        slide_count=len(prs.slides),
        file_size_mb=round(file_size, 2),
    )

    for slide_num, slide in enumerate(prs.slides, 1):
        info = SlideInfo(number=slide_num)

        # Layout name
        if slide.slide_layout:
            info.layout_name = slide.slide_layout.name

        # Title
        if slide.shapes.title:
            info.title = slide.shapes.title.text.strip()
            info.has_title = bool(info.title)

        # Process all shapes
        texts = []
        for shape in slide.shapes:
            info.shape_count += 1
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                info.image_count += 1
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)

        info.text_content = "\n".join(texts)
        info.word_count = len(info.text_content.split())

        # Notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            info.notes = slide.notes_slide.notes_text_frame.text.strip()

        inventory.slides.append(info)
        inventory.total_images += info.image_count
        inventory.total_words += info.word_count

    if output_path:
        os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
        with open(output_path, "w", encoding="utf-8") as f:
            json.dump(asdict(inventory), f, indent=2, ensure_ascii=False)
        print(f"Inventory saved to: {output_path}")

    return inventory


# ---------------------------------------------------------------------------
# Full extraction pipeline
# ---------------------------------------------------------------------------

def full_extract(pptx_path: str, output_base: str, render: bool = True,
                 width: int = 1920, height: int = 1080) -> dict:
    """Run the full extraction pipeline on a single PPTX.

    Args:
        pptx_path: Path to the PPTX file.
        output_base: Base directory for outputs (e.g., slides/01-introduction/extracted/).
        render: Whether to render slides as PNG (requires PowerPoint).
        width: Render width.
        height: Render height.

    Returns:
        Dict with paths to all generated artifacts.
    """
    os.makedirs(output_base, exist_ok=True)
    results = {"pptx": pptx_path}

    # 1. Render slides as PNG
    if render:
        renders_dir = os.path.join(output_base, "renders")
        try:
            count = render_slides(pptx_path, renders_dir, width, height)
            results["renders"] = renders_dir
            results["render_count"] = count
        except Exception as e:
            print(f"WARNING: Rendering failed: {e}")
            results["renders"] = None

    # 2. Extract text
    text_path = os.path.join(output_base, "content.md")
    extract_text(pptx_path, text_path)
    results["text"] = text_path

    # 3. Extract images
    images_dir = os.path.join(output_base, "images")
    img_count = extract_images(pptx_path, images_dir)
    results["images"] = images_dir
    results["image_count"] = img_count

    # 4. Generate inventory
    inventory_path = os.path.join(output_base, "inventory.json")
    inventory = generate_inventory(pptx_path, inventory_path)
    results["inventory"] = inventory_path
    results["slide_count"] = inventory.slide_count

    return results


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(
        description="Slide tools: render, extract, and analyze PowerPoint presentations.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    subparsers = parser.add_subparsers(dest="command", help="Command to execute")

    # render
    p_render = subparsers.add_parser("render", help="Render slides to PNG via PowerPoint COM")
    p_render.add_argument("pptx_path", help="Path to PPTX file")
    p_render.add_argument("--outdir", help="Output directory", default=None)
    p_render.add_argument("--width", type=int, default=1920, help="Image width (default: 1920)")
    p_render.add_argument("--height", type=int, default=1080, help="Image height (default: 1080)")

    # extract
    p_extract = subparsers.add_parser("extract", help="Full extraction (render + text + images + inventory)")
    p_extract.add_argument("pptx_path", help="Path to PPTX file")
    p_extract.add_argument("--outdir", help="Output base directory", default=None)
    p_extract.add_argument("--no-render", action="store_true", help="Skip PNG rendering")

    # inventory
    p_inv = subparsers.add_parser("inventory", help="Generate JSON inventory of a PPTX")
    p_inv.add_argument("pptx_path", help="Path to PPTX file")
    p_inv.add_argument("--output", help="Output JSON path", default=None)

    # batch-render
    p_batch = subparsers.add_parser("batch-render", help="Render all PPTX in a directory")
    p_batch.add_argument("directory", help="Directory containing PPTX files")
    p_batch.add_argument("--outdir", help="Base output directory", default=None)
    p_batch.add_argument("--width", type=int, default=1920)
    p_batch.add_argument("--height", type=int, default=1080)

    # batch-extract
    p_batchex = subparsers.add_parser("batch-extract", help="Full extraction on all PPTX in a directory")
    p_batchex.add_argument("directory", help="Directory containing PPTX files")
    p_batchex.add_argument("--outdir", help="Base output directory", default=None)
    p_batchex.add_argument("--no-render", action="store_true", help="Skip PNG rendering")

    # marp-render
    p_marp = subparsers.add_parser("marp-render", help="Render Marp slides.md to PNG")
    p_marp.add_argument("deck_path", help="Path to deck directory (e.g., slides/01-introduction)")
    p_marp.add_argument("--scale", type=int, default=2, help="Image scale factor (default: 2)")

    # marp-render-all
    p_marp_all = subparsers.add_parser("marp-render-all", help="Render all Marp decks to PNG")
    p_marp_all.add_argument("--scale", type=int, default=2, help="Image scale factor (default: 2)")

    # check-refs
    p_check = subparsers.add_parser("check-refs", help="Check image references in slides.md")
    p_check.add_argument("deck_path", help="Path to deck directory")

    # compare
    p_compare = subparsers.add_parser("compare", help="Compare PPTX vs Marp slide counts")
    p_compare.add_argument("deck_path", nargs="?", help="Path to deck directory (or all if omitted)")

    args = parser.parse_args()

    if not args.command:
        parser.print_help()
        return

    if args.command == "render":
        outdir = args.outdir or os.path.join(os.path.dirname(args.pptx_path), "renders")
        render_slides(args.pptx_path, outdir, args.width, args.height)

    elif args.command == "extract":
        outdir = args.outdir or os.path.join(os.path.dirname(args.pptx_path), "extracted")
        full_extract(args.pptx_path, outdir, render=not args.no_render)

    elif args.command == "inventory":
        output = args.output or args.pptx_path.replace(".pptx", "_inventory.json")
        inv = generate_inventory(args.pptx_path, output)
        print(f"Slides: {inv.slide_count}, Images: {inv.total_images}, Words: {inv.total_words}")

    elif args.command == "batch-render":
        results = batch_render(args.directory, args.outdir, args.width, args.height)
        print(f"\nBatch complete: {len(results)} files processed")
        for name, count in results.items():
            status = f"{count} slides" if count >= 0 else "FAILED"
            print(f"  {name}: {status}")

    elif args.command == "batch-extract":
        pptx_files = sorted(Path(args.directory).rglob("*.pptx"))
        for pptx_path in pptx_files:
            if args.outdir:
                outdir = os.path.join(args.outdir, pptx_path.stem)
            else:
                outdir = os.path.join(os.path.dirname(str(pptx_path)), "extracted")
            print(f"\n--- Processing: {pptx_path.name} ---")
            full_extract(str(pptx_path), outdir, render=not args.no_render)

    elif args.command == "marp-render":
        render_marp(args.deck_path, args.scale)

    elif args.command == "marp-render-all":
        results = render_all_marp(args.scale)
        print(f"\nBatch Marp render complete:")
        for name, count in results.items():
            print(f"  {name}: {count} slides")

    elif args.command == "check-refs":
        check_image_refs(args.deck_path)

    elif args.command == "compare":
        if hasattr(args, "deck_path") and args.deck_path:
            compare_slide_counts(args.deck_path)
        else:
            print("Comparing all decks:")
            for deck_name in DECK_NAMES:
                deck_path = SLIDES_ROOT / deck_name
                if (deck_path / "slides.md").exists():
                    compare_slide_counts(str(deck_path))


if __name__ == "__main__":
    main()
